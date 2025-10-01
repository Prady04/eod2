import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import sys
from datetime import datetime

def add_footer(slide, date_text):
    """Add footer with date to slide."""
    # Footer text box
    left = Inches(0.5)
    top = Inches(7)
    width = Inches(9)
    height = Inches(0.4)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = date_text
    
    # Format footer text - Changed to light gray for black background
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(180, 180, 180)

def add_watermark(slide, text="Prady ©"):
    """Add diagonal watermark across slide (behind content)."""
    # Watermark positioned diagonally across the slide
    left = Inches(2)
    top = Inches(2.5)
    width = Inches(6)
    height = Inches(2)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    tf.word_wrap = False
    
    # Format watermark - Darker for black background
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 40, 40)
    
    # Make background transparent
    txBox.fill.background()
    txBox.line.fill.background()
    
    # Rotate the text box for diagonal effect
    txBox.rotation = 315  # -45 degrees (diagonal)
    
    # Send to back so it appears behind the image
    slide.shapes._spTree.remove(txBox._element)
    slide.shapes._spTree.insert(2, txBox._element)

def create_index_slide(prs, png_files, date_text):
    """Create an attractive index slide with image count and summary."""
    blank_layout = prs.slide_layouts[6]  # Blank layout for custom design
    slide = prs.slides.add_slide(blank_layout)
    
    # Black background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    
    # Decorative circles in corners
    # Top-left circle
    circle1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-0.5), Inches(-0.5), Inches(1.5), Inches(1.5)
    )
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RGBColor(70, 130, 255)
    circle1.fill.transparency = 0.3
    circle1.line.fill.background()
    
    # Bottom-right circle
    circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9), Inches(6.5), Inches(1.5), Inches(1.5)
    )
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(255, 100, 100)
    circle2.fill.transparency = 0.3
    circle2.line.fill.background()
    
    # Top-right small circle
    circle3 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(8.5), Inches(0.5), Inches(1), Inches(1)
    )
    circle3.fill.solid()
    circle3.fill.fore_color.rgb = RGBColor(100, 255, 150)
    circle3.fill.transparency = 0.4
    circle3.line.fill.background()
    
    # Main Title
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.text = "Image Collection"
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Description text
    left = Inches(1.5)
    top = Inches(3.2)
    width = Inches(7)
    height = Inches(0.8)
    
    desc_box = slide.shapes.add_textbox(left, top, width, height)
    tf = desc_box.text_frame
    tf.text = "Monthly breakdown of positive days with momentum scores"
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(180, 220, 255)
    
    # Subtitle with count
    left = Inches(1)
    top = Inches(4.3)
    width = Inches(8)
    height = Inches(1)
    
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    tf = subtitle_box.text_frame
    tf.text = f"{len(png_files)} Images"
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 200, 255)
    
    # Decorative line with gradient effect (using rectangles)
    left = Inches(2.5)
    top = Inches(5.3)
    width = Inches(5)
    height = Inches(0.08)
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(70, 180, 255)
    line.line.fill.background()
    
    # Add decorative squares
    square_size = Inches(0.3)
    colors = [
        RGBColor(70, 180, 255),
        RGBColor(100, 200, 255),
        RGBColor(150, 220, 255)
    ]
    
    for i, color in enumerate(colors):
        square = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(3.5 + i * 0.8), Inches(5.8), square_size, square_size
        )
        square.fill.solid()
        square.fill.fore_color.rgb = color
        square.line.fill.background()
    
    # Add footer to index
    add_footer(slide, date_text)
    
    print("Created index slide with graphics")

def create_ppt_from_images(folder_path, output_name="output.pptx"):
    """
    Create a PowerPoint presentation from all PNG files in a folder.
    
    Args:
        folder_path: Path to folder containing PNG files
        output_name: Name of output PowerPoint file
    """
    # Verify folder exists
    if not os.path.exists(folder_path):
        print(f"Error: Folder '{folder_path}' does not exist!")
        return
    
    # Get all PNG files and sort them
    png_files = sorted([f for f in os.listdir(folder_path) 
                       if f.lower().endswith('.png')])
    
    if not png_files:
        print(f"No PNG files found in '{folder_path}'")
        return
    
    print(f"Found {len(png_files)} PNG files")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Get today's date
    today = datetime.now().strftime("%B %d, %Y")
    
    # Create index slide first
    create_index_slide(prs, png_files, today)
    
    # Add each image to a slide
    for idx, png_file in enumerate(png_files, 1):
        # Create blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Black background for all slides
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)
        
        # Full path to image
        img_path = os.path.join(folder_path, png_file)
        
        # Add image centered on slide
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(6)
        
        try:
            slide.shapes.add_picture(img_path, left, top, width=width, height=height)
            
            # Add watermark
            add_watermark(slide, "Prady ©")
            
            # Add footer with date
            add_footer(slide, today)
            
            print(f"Added slide {idx + 1}/{len(png_files) + 1}: {png_file}")
        except Exception as e:
            print(f"Error adding {png_file}: {e}")
    
    # Save presentation
    output_path = os.path.join(folder_path, output_name)
    prs.save(output_path)
    print(f"\nPresentation saved as: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python script.py <folder_path> [output_name.pptx]")
        print("Example: python script.py ./images my_presentation.pptx")
        sys.exit(1)
    
    folder = sys.argv[1]
    output = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"
    
    create_ppt_from_images(folder, output)